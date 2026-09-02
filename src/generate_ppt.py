"""
自动生成 .pptx 答辩PPT，约19页。
使用 python-pptx，每页包含标题、要点文字和相关图表。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from config import FIGURES_DIR, SLIDES_DIR, MODEL_NAMES_CN


# 颜色方案
BG_COLOR = RGBColor(0xF8, 0xF9, 0xFA)
TITLE_COLOR = RGBColor(0x2C, 0x3E, 0x50)
ACCENT_COLOR = RGBColor(0x34, 0x98, 0xDB)
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _add_title_slide(prs, title, subtitle):
    """封面页。"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景色
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = '微软雅黑'

    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0xEC, 0xF0, 0xF1)
    run2.font.name = '微软雅黑'

    # 底部信息
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = '计算机学院\n2026年9月'
    run3.font.size = Pt(16)
    run3.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    run3.font.name = '微软雅黑'


def _add_content_slide(prs, title, bullet_points, image_path=None, image_width=Inches(6)):
    """内容页：标题 + 要点 + 可选图片。"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景色
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR

    # 标题栏
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    run.font.name = '微软雅黑'

    # 标题下划线
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.5), Inches(1.1),
        Inches(9), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()

    # 要点文字
    left = Inches(0.6)
    top = Inches(1.4)
    width = Inches(5) if image_path else Inches(9)
    height = Inches(5.5)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = point
        run2.font.size = Pt(16)
        run2.font.color.rgb = TEXT_COLOR
        run2.font.name = '微软雅黑'
        p2.space_after = Pt(8)

    # 图片
    if image_path and os.path.exists(image_path):
        pic_left = Inches(5.5)
        pic_top = Inches(1.5)
        slide.shapes.add_picture(image_path, pic_left, pic_top, width=image_width)


def _add_image_slide(prs, title, image_path, caption='', image_width=Inches(7)):
    """图片展示页：标题 + 大图。"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    run.font.name = '微软雅黑'

    # 图片
    if image_path and os.path.exists(image_path):
        pic_left = Inches(1.2)
        pic_top = Inches(1.2)
        slide.shapes.add_picture(image_path, pic_left, pic_top, width=image_width)

    # 说明
    if caption:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = caption
        run2.font.size = Pt(12)
        run2.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
        run2.font.name = '微软雅黑'


def _add_table_slide(prs, title, metrics_df):
    """表格展示页。"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    run.font.name = '微软雅黑'

    # 表格
    rows = len(metrics_df) + 1
    cols = len(metrics_df.columns)
    table_left = Inches(0.8)
    table_top = Inches(1.5)
    table_width = Inches(8.4)
    table_height = Inches(3.5)
    table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

    # 表头
    for j, col in enumerate(metrics_df.columns):
        cell = table.cell(0, j)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = '微软雅黑'

    # 数据行
    for i, row in metrics_df.iterrows():
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEC, 0xF0, 0xF1)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(13)
                    run.font.color.rgb = TEXT_COLOR
                    run.font.name = '微软雅黑'


def generate_ppt(metrics_all, metrics_df, feature_names, target_names):
    """生成答辩PPT。"""
    print("[PPT生成] 开始生成答辩PPT...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 页1: 封面
    _add_title_slide(prs,
                     '鸢尾花分类——从模型对比到可解释性',
                     '机器学习项目实训课程设计答辩')

    # 页2: 目录
    _add_content_slide(prs, '目录', [
        '一、项目背景与数据集介绍',
        '二、数据探索与预处理',
        '三、模型原理概述',
        '四、模型训练与评估',
        '五、决策边界可视化对比',
        '六、可解释性分析（决策树规则 + SHAP）',
        '七、结论与展望',
    ])

    # 页3: 项目背景
    _add_content_slide(prs, '一、项目背景', [
        '鸢尾花数据集（Iris）由Fisher于1936年引入',
        '机器学习领域最经典的入门基准数据集',
        '150条样本，4个特征，3个类别',
        '本课题核心目标：',
        '  • 系统对比四种分类器的决策边界',
        '  • 提取可理解的决策树分类规则',
        '  • 利用SHAP进行可解释性分析',
        '  • 不仅看准确率，更关注可解释性',
    ])

    # 页4: 数据集介绍
    _add_content_slide(prs, '二、数据集介绍', [
        '特征（4个连续型，单位cm）：',
        '  • 花萼长度  花萼宽度',
        '  • 花瓣长度  花瓣宽度',
        '',
        '类别（3类，各50条，均衡分布）：',
        '  • 山鸢尾 (Setosa)',
        '  • 变色鸢尾 (Versicolor)',
        '  • 维吉尼亚鸢尾 (Virginica)',
        '',
        '数据来源：UCI / sklearn内置',
        '划分策略：7:3分层抽样（训练105/测试45）',
    ])

    # 页5: EDA-散点矩阵
    _add_image_slide(prs, '探索性数据分析——散点矩阵',
                     os.path.join(FIGURES_DIR, '03_scatter_matrix.png'),
                     '山鸢尾完全线性可分，变色与维吉尼亚在花瓣特征空间分离度更好',
                     image_width=Inches(7.5))

    # 页6: EDA-相关性与箱线
    _add_image_slide(prs, '探索性数据分析——相关性热力图',
                     os.path.join(FIGURES_DIR, '04_correlation_heatmap.png'),
                     '花瓣长度与花瓣宽度强正相关(r=0.96)',
                     image_width=Inches(6))

    # 页7: 预处理策略
    _add_content_slide(prs, '三、数据预处理策略', [
        '1. 训练/测试集划分',
        '  • 分层抽样 stratify=y, test_size=0.3',
        '  • 训练集105条, 测试集45条, 各类均衡',
        '',
        '2. 特征标准化 (StandardScaler)',
        '  • KNN/SVM/逻辑回归: 使用标准化数据',
        '  • 决策树: 使用原始数据 (保持规则阈值可读性)',
        '  • 仅fit训练集, transform测试集',
        '',
        '3. 随机种子 random_state=42, 结果可复现',
    ])

    # 页8: 四模型原理
    _add_content_slide(prs, '四、模型原理速览', [
        'KNN (K=5): 基于实例的懒惰学习, 多数投票',
        '  → 边界碎片化, 适应任意分布',
        '',
        'SVM-RBF (C=1.0, gamma=scale): 核函数映射高维空间',
        '  → 边界平滑, 泛化能力通常最优',
        '',
        '决策树 (max_depth=3, Gini): 递归特征分裂',
        '  → 边界阶梯状, 规则可读',
        '',
        '逻辑回归 (Softmax, lbfgs): 线性分类器',
        '  → 边界为直线, 简单高效',
    ])

    # 页9: 性能对比表
    _add_table_slide(prs, '五、模型性能对比', metrics_df)

    # 页10: 混淆矩阵
    _add_image_slide(prs, '混淆矩阵分析',
                     os.path.join(FIGURES_DIR, '05_confusion_matrices.png'),
                     '山鸢尾100%正确, 错误集中在变色与维吉尼亚之间',
                     image_width=Inches(7.5))

    # 页11: ROC曲线
    _add_image_slide(prs, 'ROC曲线对比',
                     os.path.join(FIGURES_DIR, '06_roc_curves.png'),
                     '所有模型AUC>0.99, SVM和决策树AUC最高(0.9951)',
                     image_width=Inches(6))

    # 页12: 决策边界-花瓣
    _add_image_slide(prs, '决策边界——花瓣特征对（高区分度）',
                     os.path.join(FIGURES_DIR, '08_decision_boundary_高区分度.png'),
                     'KNN碎片化 / SVM平滑 / 决策树阶梯 / 逻辑回归线性',
                     image_width=Inches(7.5))

    # 页13: 决策边界-花萼
    _add_image_slide(prs, '决策边界——花萼特征对（低区分度）',
                     os.path.join(FIGURES_DIR, '09_decision_boundary_低区分度.png'),
                     '花萼特征区分度低, 各模型边界混乱',
                     image_width=Inches(7.5))

    # 页14: 边界差异总结
    _add_content_slide(prs, '五、边界形态差异总结', [
        'KNN: 碎片化边界',
        '  • 由局部邻居投票决定, 对噪声敏感',
        '',
        'SVM: 平滑非线性边界',
        '  • RBF核映射高维空间, 泛化性最好',
        '',
        '决策树: 阶梯状边界',
        '  • 轴平行线段, 规则可读但无法斜线划分',
        '',
        '逻辑回归: 线性边界',
        '  • 超平面划分, 仅适合线性可分场景',
    ])

    # 页15: 决策树规则
    _add_image_slide(prs, '六、决策树规则提取',
                     os.path.join(FIGURES_DIR, '10_tree_structure.png'),
                     '花瓣长度≤2.45cm→Setosa | 花瓣宽度≤1.75cm→Versicolor | 否则→Virginica',
                     image_width=Inches(7.5))

    # 页16: SHAP全局
    _add_image_slide(prs, 'SHAP全局解释——SVM',
                     os.path.join(FIGURES_DIR, '12_shap_summary_svm.png'),
                     '花瓣长度和花瓣宽度的SHAP值远大于花萼特征',
                     image_width=Inches(6))

    # 页17: SHAP局部
    _add_image_slide(prs, 'SHAP局部解释——单样本决策',
                     os.path.join(FIGURES_DIR, '16_shap_force_plot.png'),
                     '各特征SHAP值逐步将预测从基础值推向最终类别',
                     image_width=Inches(7))

    # 页18: 结论与展望
    _add_content_slide(prs, '七、结论与展望', [
        '实验结论：',
        '  • 决策树和SVM性能最优, SVM泛化能力最强',
        '  • 花瓣特征是分类的关键（EDA/树重要性/SHAP一致确认）',
        '  • 决策树规则+SHAP实现多层面可解释性',
        '',
        '局限性：',
        '  • 数据集规模小(150条), 统计差异不显著',
        '  • 未进行超参数调优',
        '',
        '改进方向：',
        '  • 引入交叉验证和网格搜索',
        '  • 增加集成学习模型对比',
        '  • 探索LIME等其他可解释性方法',
    ])

    # 页19: 谢谢
    _add_title_slide(prs, '谢谢', '欢迎提问与指导')

    # 保存
    path = os.path.join(SLIDES_DIR, '鸢尾花分类答辩PPT.pptx')
    prs.save(path)
    print(f"  -> {path}")
